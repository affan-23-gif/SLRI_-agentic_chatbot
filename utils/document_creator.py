import os
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from docx import Document as DocxDocument # Avoid conflict with python-docx Document
from pptx import Presentation # pip install python-pptx

def create_dummy_documents(base_dir: str = "test_docs"):
    """
    Creates various dummy document files for testing the IngestionAgent.
    """
    os.makedirs(base_dir, exist_ok=True)
    print(f"Creating dummy documents in: {os.path.abspath(base_dir)}")

    # 1. TXT File
    txt_path = os.path.join(base_dir, "example.txt")
    with open(txt_path, "w") as f:
        f.write("This is a simple text document. It contains information about Q1 KPIs. Revenue was up by 10%. Customer retention is a key focus.")
    print(f"Generated {txt_path}")

    # 2. Markdown File
    md_path = os.path.join(base_dir, "sales_review.md")
    with open(md_path, "w") as f:
        f.write("# Sales Review Q1\n\n"
                  "This document summarizes the sales performance for the first quarter.\n\n"
                  "* **Revenue**: $1.2 Million\n"
                  "* **Customer Acquisition Cost (CAC)**: $50 per customer\n"
                  "* **Conversion Rate**: 5%\n\n"
                  "These are the key performance indicators for the first quarter. We also tracked Net Promoter Score (NPS).")
    print(f"Generated {md_path}")

    # 3. PDF File (using ReportLab)
    pdf_path = os.path.join(base_dir, "metrics.pdf")
    try:
        c = canvas.Canvas(pdf_path, pagesize=letter)
        c.drawString(100, 750, "Quarterly Metrics Report - Q1")
        c.drawString(100, 730, "Key Performance Indicators (KPIs):")
        c.drawString(120, 710, "- Revenue: $1.5 Million")
        c.drawString(120, 690, "- Customer Retention: 85%")
        c.drawString(120, 670, "- Net Promoter Score (NPS): 45")
        c.drawString(100, 600, "This report provides an overview of our performance.")
        c.save()
        print(f"Generated {pdf_path}")
    except ImportError:
        print("reportlab not installed. Skipping PDF generation. Please install 'reportlab' or manually place a PDF in test_docs/ if needed for testing.")
    except Exception as e:
        print(f"Error generating PDF {pdf_path}: {e}")

    # 4. DOCX File (using python-docx)
    docx_path = os.path.join(base_dir, "project_summary.docx")
    try:
        document = DocxDocument()
        document.add_heading('Project Summary - Phase 1', level=1)
        document.add_paragraph('This document outlines the key achievements and challenges of Phase 1 of the project.')
        document.add_paragraph('Key achievements include successful deployment of Module A and completion of user acceptance testing.')
        document.add_paragraph('Challenges faced were integration issues with legacy systems and resource allocation.')
        document.save(docx_path)
        print(f"Generated {docx_path}")
    except Exception as e:
        print(f"Error generating DOCX {docx_path}: {e}")

    # 5. CSV File (using pandas)
    csv_path = os.path.join(base_dir, "employee_data.csv")
    try:
        import pandas as pd
        data = {
            'EmployeeID': [1, 2, 3],
            'Name': ['Alice', 'Bob', 'Charlie'],
            'Department': ['HR', 'Engineering', 'Sales'],
            'Salary': [70000, 120000, 90000]
        }
        df = pd.DataFrame(data)
        df.to_csv(csv_path, index=False)
        print(f"Generated {csv_path}")
    except ImportError:
        print("pandas not installed. Skipping CSV generation. Please install 'pandas' or manually place a CSV in test_docs/ if needed for testing.")
    except Exception as e:
        print(f"Error generating CSV {csv_path}: {e}")

    # 6. PPTX File (using python-pptx)
    pptx_path = os.path.join(base_dir, "presentation_slides.pptx")
    try:
        prs = Presentation()
        
        # Slide 1
        slide1_layout = prs.slide_layouts[0] # Title slide
        slide1 = prs.slides.add_slide(slide1_layout)
        title1 = slide1.shapes.title
        subtitle1 = slide1.placeholders[1]
        title1.text = "Company Overview"
        subtitle1.text = "Key Highlights for the Year"

        # Slide 2
        slide2_layout = prs.slide_layouts[1] # Title and Content
        slide2 = prs.slides.add_slide(slide2_layout)
        title2 = slide2.shapes.title
        body2 = slide2.placeholders[1]
        title2.text = "Financial Performance"
        tf = body2.text_frame
        p = tf.add_paragraph()
        p.text = "Revenue growth of 20%."
        p = tf.add_paragraph()
        p.text = "Profit margins improved by 5%."

        prs.save(pptx_path)
        print(f"Generated {pptx_path}")
    except Exception as e:
        print(f"Error generating PPTX {pptx_path}: {e}")

